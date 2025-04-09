import sys
import os
from docx import Document
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from io import BytesIO
from PIL import Image

# Input DOCX path
docx_path = sys.argv[1]

# Output paths
base_name = os.path.splitext(os.path.basename(docx_path))[0]
output_dir = os.path.join(os.path.dirname(__file__), "converted")
os.makedirs(output_dir, exist_ok=True)
pptx_path = os.path.join(output_dir, f"{base_name}.pptx")

# Load Word and Presentation
doc = Document(docx_path)
prs = Presentation()

# Styles
TITLE_SIZE = Pt(24)
CONTENT_SIZE = Pt(22)
TITLE_COLOR = RGBColor(59, 130, 246)
CONTENT_COLOR = RGBColor(0, 0, 0)
NAVBAR_COLOR = RGBColor(116, 2, 254)
FONT_NAME = "Calibri"
MAX_BULLETS_PER_SLIDE = 6

# Title slide
title_slide = prs.slides.add_slide(prs.slide_layouts[0])
title_slide.shapes.title.text = "Generated Presentation"
title_slide.placeholders[1].text = os.path.basename(docx_path)

# Add navbar to slide
def add_navbar(slide):
    shape = slide.shapes.add_shape(
        autoshape_type_id=1,
        left=Inches(0),
        top=Inches(0),
        width=prs.slide_width,
        height=Inches(0.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVBAR_COLOR
    shape.line.fill.background()

    left_box = slide.shapes.add_textbox(Inches(0.2), Inches(0.05), Inches(4), Inches(0.4))
    left_tf = left_box.text_frame
    left_tf.text = "PHELP"
    p1 = left_tf.paragraphs[0]
    p1.font.size = Pt(16)
    p1.font.color.rgb = RGBColor(255, 255, 255)
    p1.font.bold = True

    right_box = slide.shapes.add_textbox(prs.slide_width - Inches(3.5), Inches(0.05), Inches(3.2), Inches(0.4))
    right_tf = right_box.text_frame
    right_tf.text = "Powered by PATEL"
    p2 = right_tf.paragraphs[0]
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(255, 255, 255)
    p2.font.bold = True

# Add slide with bullets
def add_slide(title_text, bullet_points):
    chunks = [bullet_points[i:i + MAX_BULLETS_PER_SLIDE] for i in range(0, len(bullet_points), MAX_BULLETS_PER_SLIDE)]
    for i, chunk in enumerate(chunks):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_navbar(slide)

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(9), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"{title_text}" if i == 0 else f"{title_text} (cont.)"
        run.font.size = TITLE_SIZE
        run.font.name = FONT_NAME
        run.font.bold = True
        run.font.color.rgb = TITLE_COLOR

        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.5), Inches(5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for bullet in chunk:
            p = tf.add_paragraph()
            p.level = 0
            run = p.add_run()
            run.text = f"• {bullet}"
            run.font.size = CONTENT_SIZE
            run.font.name = FONT_NAME
            run.font.color.rgb = CONTENT_COLOR

# Add image slide
def add_image_slide(image_blob, title_text="Image"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_navbar(slide)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1))
    tf = title_box.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = title_text
    run.font.size = TITLE_SIZE
    run.font.name = FONT_NAME
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR

    # Save image temporarily
    image = Image.open(BytesIO(image_blob))
    temp_image_path = os.path.join(output_dir, "temp_image.jpg")
    image.save(temp_image_path)

    # Insert image
    slide.shapes.add_picture(temp_image_path, Inches(1), Inches(1.7), width=Inches(7.5))

    os.remove(temp_image_path)

# DOCX processing
current_title = "Untitled Slide"
bullets = []

# Detect images via document relationships
rels = doc.part._rels
image_parts = [rel.target_part for rel in rels.values() if "image" in rel.reltype]

image_idx = 0

for para in doc.paragraphs:
    text = para.text.strip()
    if not text:
        continue

    is_heading = para.style.name.startswith("Heading")
    is_bold = any(run.bold for run in para.runs)

    if is_heading or is_bold:
        if bullets:
            add_slide(current_title, bullets)
            bullets = []
        current_title = text
    else:
        bullets.append(text)

# Add any remaining bullets
if bullets:
    add_slide(current_title, bullets)

# Add each image as a separate slide
for part in image_parts:
    image_blob = part.blob
    add_image_slide(image_blob, title_text=f"Image {image_idx + 1}")
    image_idx += 1

# Save final PPTX
prs.save(pptx_path)
print(f"Perfectly styled PPT with images saved to: {pptx_path}")
